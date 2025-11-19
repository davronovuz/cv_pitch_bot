import os
import asyncio
import sqlite3
import logging
import json
from datetime import datetime
from typing import Optional, Dict, List
from aiogram.types import InlineKeyboardMarkup, InlineKeyboardButton, InputFile, ContentType
from aiogram.contrib.fsm_storage.memory import MemoryStorage
from aiogram.dispatcher import FSMContext
from aiogram.dispatcher.filters.state import State, StatesGroup
from loader import bot, dp
from aiogram import types

from environs import Env

# environs kutubxonasidan foydalanish
env = Env()
env.read_env()

# .env fayl ichidan quyidagilarni o'qiymiz
OPENAI_API_KEY = env.str("OPENAI_API_KEY")  # Bot toekn

try:
    from openai import OpenAI
except ImportError:
    OpenAI = None

# ==================== KONFIGURATSIYA ====================
ADMIN_ID = 1879114908
USE_OPENAI = True  # Yoqilgan

CURRENCY = "so'm"
MAX_FILE_SIZE = 10 * 1024 * 1024  # 10 MB
ALLOWED_FILE_TYPES = ['image/jpeg', 'image/png', 'application/pdf']

# Karta ma'lumotlari
CARD_NUMBER = "8600 1234 5678 9012"
CARD_HOLDER = "JOHN DOE"

# Logging sozlash
logging.basicConfig(level=logging.INFO)
logger = logging.getLogger(__name__)

# OpenAI client
openai_client = None
if USE_OPENAI and OpenAI and OPENAI_API_KEY:
    openai_client = OpenAI(api_key=OPENAI_API_KEY)

# ==================== SAVOLLAR RO'YXATI ====================
QUESTIONS = [
    "1️⃣ Ismingiz?",
    "2️⃣ Loyiha nomi?",
    "3️⃣ Loyiha tavsifi (qisqacha, 2-3 jumla)?",
    "4️⃣ Qanday muammoni hal qilasiz?",
    "5️⃣ Sizning yechimingiz?",
    "6️⃣ Maqsadli auditoriya kimlar?",
    "7️⃣ Biznes model (qanday daromad olasiz)?",
    "8️⃣ Asosiy raqobatchilaringiz?",
    "9️⃣ Sizning ustunligingiz (raqobatchilardan farqi)?",
    "🔟 Moliyaviy prognoz (keyingi 1 yil)?",
]


# [Database class o'zgarmaydi - yuqoridagi koddan olish kerak]

# ==================== AI ENHANCEMENT FUNKSIYALARI ====================

async def enhance_answers_with_openai(answers: List[str], package: str) -> List[str]:
    """Javoblarni OpenAI bilan professional qilish"""
    if not USE_OPENAI or not openai_client:
        return answers

    # Pro paket uchun kuchliroq model
    model = "gpt-4" if package == "pro" else "gpt-3.5-turbo"

    try:
        # Har bir javobni yaxshilash
        enhanced_answers = []

        prompts = [
            f"Make this founder name more professional (keep it short): {answers[0]}",
            f"Create a catchy and professional project name based on: {answers[1]}",
            f"Transform this project description into a compelling elevator pitch (2-3 sentences max): {answers[2]}",
            f"Rewrite this problem statement to make it more urgent and clear for investors (3-4 bullet points): {answers[3]}",
            f"Transform this solution into clear value propositions (3-4 bullet points): {answers[4]}",
            f"Define this target audience with specific demographics and market size: {answers[5]}",
            f"Create a clear revenue model explanation with potential revenue streams: {answers[6]}",
            f"Analyze these competitors and create a competitive landscape overview: {answers[7]}",
            f"Transform these advantages into unique selling propositions (USPs): {answers[8]}",
            f"Create realistic financial projections with key metrics for next year: {answers[9]}"
        ]

        for i, prompt in enumerate(prompts):
            if i < len(answers):
                response = await asyncio.to_thread(
                    lambda: openai_client.chat.completions.create(
                        model=model,
                        messages=[
                            {"role": "system",
                             "content": "You are a professional pitch deck consultant. Keep responses concise and impactful."},
                            {"role": "user", "content": prompt}
                        ],
                        max_tokens=200,
                        temperature=0.7
                    )
                )
                enhanced_answers.append(response.choices[0].message.content.strip())
            else:
                enhanced_answers.append(answers[i] if i < len(answers) else "")

        return enhanced_answers

    except Exception as e:
        logger.error(f"OpenAI enhancement failed: {e}")
        return answers


async def create_professional_pitch_content(answers: List[str], package: str) -> Dict:
    """AI bilan to'liq pitch content yaratish"""

    if not USE_OPENAI or not openai_client:
        return {
            'project_name': answers[1] if len(answers) > 1 else "Startup",
            'author': answers[0] if len(answers) > 0 else "Entrepreneur",
            'tagline': answers[2][:100] if len(answers) > 2 else "Innovative Solution",
            'problem': answers[3] if len(answers) > 3 else "",
            'solution': answers[4] if len(answers) > 4 else "",
            'market': answers[5] if len(answers) > 5 else "",
            'business_model': answers[6] if len(answers) > 6 else "",
            'competition': answers[7] if len(answers) > 7 else "",
            'advantage': answers[8] if len(answers) > 8 else "",
            'financials': answers[9] if len(answers) > 9 else "",
            'cta': "Let's build the future together! 🚀",
        }

    model = "gpt-4" if package == "pro" else "gpt-3.5-turbo"

    prompt = f"""
Create a professional investor pitch based on these startup details. 
Make it compelling, data-driven, and investor-ready.

Founder: {answers[0] if len(answers) > 0 else ""}
Project: {answers[1] if len(answers) > 1 else ""}
Description: {answers[2] if len(answers) > 2 else ""}
Problem: {answers[3] if len(answers) > 3 else ""}
Solution: {answers[4] if len(answers) > 4 else ""}
Target Market: {answers[5] if len(answers) > 5 else ""}
Business Model: {answers[6] if len(answers) > 6 else ""}
Competition: {answers[7] if len(answers) > 7 else ""}
Our Advantage: {answers[8] if len(answers) > 8 else ""}
Financial Forecast: {answers[9] if len(answers) > 9 else ""}

Return a JSON object with these exact keys:
{{
  "project_name": "compelling project name",
  "author": "founder name",
  "tagline": "powerful tagline (max 10 words)",
  "problem": "3-4 bullet points about the problem",
  "solution": "3-4 bullet points about the solution",
  "market": "target market analysis with TAM/SAM/SOM",
  "business_model": "revenue streams and pricing strategy",
  "competition": "competitive analysis",
  "advantage": "3 key differentiators",
  "financials": "key metrics and projections",
  "cta": "powerful call to action"
}}

Make it professional and investor-ready. Use bullet points where indicated.
"""

    try:
        response = await asyncio.to_thread(
            lambda: openai_client.chat.completions.create(
                model=model,
                messages=[
                    {"role": "system",
                     "content": "You are a top-tier pitch deck consultant. Create compelling, professional content."},
                    {"role": "user", "content": prompt}
                ],
                max_tokens=1500,
                temperature=0.7,
                response_format={"type": "json_object"}  # JSON formatni majburlash
            )
        )

        content = json.loads(response.choices[0].message.content)
        return content

    except Exception as e:
        logger.error(f"AI content creation failed: {e}")
        # Fallback
        return {
            'project_name': answers[1] if len(answers) > 1 else "Startup",
            'author': answers[0] if len(answers) > 0 else "Entrepreneur",
            'tagline': answers[2][:100] if len(answers) > 2 else "Innovative Solution",
            'problem': answers[3] if len(answers) > 3 else "",
            'solution': answers[4] if len(answers) > 4 else "",
            'market': answers[5] if len(answers) > 5 else "",
            'business_model': answers[6] if len(answers) > 6 else "",
            'competition': answers[7] if len(answers) > 7 else "",
            'advantage': answers[8] if len(answers) > 8 else "",
            'financials': answers[9] if len(answers) > 9 else "",
            'cta': "Let's build the future together! 🚀",
        }


# ==================== YANGILANGAN PPTX YARATISH ====================
async def create_stunning_pitch_deck(user_id: int, answers: List[str], package: str) -> str:
    """Mukammal PPTX yaratish - AI optimized"""
    from pptx import Presentation
    from pptx.util import Inches, Pt
    from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
    from pptx.dml.color import RGBColor
    from pptx.enum.shapes import MSO_SHAPE
    from pptx.enum.dml import MSO_THEME_COLOR

    # AI orqali content olish
    content = await create_professional_pitch_content(answers, package)

    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)

    # Professional ranglar
    COLORS = {
        'primary': RGBColor(46, 64, 83),  # Dark blue-gray
        'secondary': RGBColor(52, 152, 219),  # Bright blue
        'accent': RGBColor(46, 204, 113),  # Green
        'danger': RGBColor(231, 76, 60),  # Red
        'warning': RGBColor(243, 156, 18),  # Orange
        'dark': RGBColor(44, 62, 80),  # Dark
        'light': RGBColor(236, 240, 241),  # Light gray
        'white': RGBColor(255, 255, 255)  # White
    }

    def add_gradient_background(slide, color1, color2):
        """Gradient fon qo'shish"""
        fill = slide.background.fill
        fill.gradient()
        fill.gradient_angle = 135
        stops = fill.gradient_stops
        stops[0].color.rgb = color1
        stops[0].position = 0.0
        stops[1].color.rgb = color2
        stops[1].position = 1.0

    # ==================== 1. TITLE SLIDE ====================
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_gradient_background(slide, COLORS['primary'], COLORS['secondary'])

    # Decorative shape
    shape = slide.shapes.add_shape(
        MSO_SHAPE.HEXAGON,
        Inches(7.5), Inches(0.5),
        Inches(2), Inches(2)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = COLORS['accent']
    shape.fill.transparency = 0.7
    shape.line.fill.background()

    # Title
    title_box = slide.shapes.add_textbox(Inches(1), Inches(2), Inches(8), Inches(1.5))
    tf = title_box.text_frame
    tf.text = content['project_name'].upper()
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    p.font.name = "Arial Black"
    p.font.size = Pt(48)
    p.font.color.rgb = COLORS['white']
    p.font.bold = True

    # Tagline
    tagline_box = slide.shapes.add_textbox(Inches(1.5), Inches(3.8), Inches(7), Inches(0.8))
    tf = tagline_box.text_frame
    tf.text = content['tagline']
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    p.font.size = Pt(24)
    p.font.color.rgb = COLORS['light']
    p.font.italic = True

    # Author
    author_box = slide.shapes.add_textbox(Inches(1), Inches(6), Inches(8), Inches(0.5))
    tf = author_box.text_frame
    tf.text = f"Presented by {content['author']}"
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    p.font.size = Pt(18)
    p.font.color.rgb = COLORS['light']

    # ==================== 2. PROBLEM SLIDE ====================
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = COLORS['white']

    # Header bar
    header = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(0), Inches(0),
        Inches(10), Inches(1.2)
    )
    header.fill.solid()
    header.fill.fore_color.rgb = COLORS['danger']
    header.line.fill.background()

    # Title
    title_box = slide.shapes.add_textbox(Inches(1), Inches(0.3), Inches(8), Inches(0.7))
    tf = title_box.text_frame
    tf.text = "🔥 THE PROBLEM"
    p = tf.paragraphs[0]
    p.font.size = Pt(36)
    p.font.color.rgb = COLORS['white']
    p.font.bold = True

    # Problem content
    content_box = slide.shapes.add_textbox(Inches(1), Inches(2), Inches(8), Inches(4.5))
    tf = content_box.text_frame
    tf.text = content['problem']
    tf.word_wrap = True
    for p in tf.paragraphs:
        p.font.size = Pt(20)
        p.font.color.rgb = COLORS['dark']
        p.space_before = Pt(12)
        p.space_after = Pt(12)

    # ==================== 3. SOLUTION SLIDE ====================
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = COLORS['white']

    # Header bar
    header = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(0), Inches(0),
        Inches(10), Inches(1.2)
    )
    header.fill.solid()
    header.fill.fore_color.rgb = COLORS['accent']
    header.line.fill.background()

    # Title
    title_box = slide.shapes.add_textbox(Inches(1), Inches(0.3), Inches(8), Inches(0.7))
    tf = title_box.text_frame
    tf.text = "💡 OUR SOLUTION"
    p = tf.paragraphs[0]
    p.font.size = Pt(36)
    p.font.color.rgb = COLORS['white']
    p.font.bold = True

    # Solution content
    content_box = slide.shapes.add_textbox(Inches(1), Inches(2), Inches(8), Inches(4.5))
    tf = content_box.text_frame
    tf.text = content['solution']
    tf.word_wrap = True
    for p in tf.paragraphs:
        p.font.size = Pt(20)
        p.font.color.rgb = COLORS['dark']
        p.space_before = Pt(12)
        p.space_after = Pt(12)

    # ==================== 4. MARKET SLIDE ====================
    if content.get('market'):
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        slide.background.fill.solid()
        slide.background.fill.fore_color.rgb = COLORS['white']

        # Header
        header = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(0), Inches(0),
            Inches(10), Inches(1.2)
        )
        header.fill.solid()
        header.fill.fore_color.rgb = COLORS['secondary']
        header.line.fill.background()

        # Title
        title_box = slide.shapes.add_textbox(Inches(1), Inches(0.3), Inches(8), Inches(0.7))
        tf = title_box.text_frame
        tf.text = "🎯 TARGET MARKET"
        p = tf.paragraphs[0]
        p.font.size = Pt(36)
        p.font.color.rgb = COLORS['white']
        p.font.bold = True

        # Market content
        content_box = slide.shapes.add_textbox(Inches(1), Inches(2), Inches(8), Inches(4.5))
        tf = content_box.text_frame
        tf.text = content['market']
        tf.word_wrap = True
        for p in tf.paragraphs:
            p.font.size = Pt(20)
            p.font.color.rgb = COLORS['dark']
            p.space_before = Pt(12)

    # ==================== 5. BUSINESS MODEL ====================
    if content.get('business_model'):
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        slide.background.fill.solid()
        slide.background.fill.fore_color.rgb = COLORS['white']

        # Header
        header = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(0), Inches(0),
            Inches(10), Inches(1.2)
        )
        header.fill.solid()
        header.fill.fore_color.rgb = COLORS['warning']
        header.line.fill.background()

        # Title
        title_box = slide.shapes.add_textbox(Inches(1), Inches(0.3), Inches(8), Inches(0.7))
        tf = title_box.text_frame
        tf.text = "💼 BUSINESS MODEL"
        p = tf.paragraphs[0]
        p.font.size = Pt(36)
        p.font.color.rgb = COLORS['white']
        p.font.bold = True

        # Business model content
        content_box = slide.shapes.add_textbox(Inches(1), Inches(2), Inches(8), Inches(4.5))
        tf = content_box.text_frame
        tf.text = content['business_model']
        tf.word_wrap = True
        for p in tf.paragraphs:
            p.font.size = Pt(20)
            p.font.color.rgb = COLORS['dark']
            p.space_before = Pt(12)

    # ==================== 6. COMPETITIVE ADVANTAGE ====================
    if content.get('advantage'):
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        slide.background.fill.solid()
        slide.background.fill.fore_color.rgb = COLORS['white']

        # Header
        header = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(0), Inches(0),
            Inches(10), Inches(1.2)
        )
        header.fill.solid()
        header.fill.fore_color.rgb = COLORS['primary']
        header.line.fill.background()

        # Title
        title_box = slide.shapes.add_textbox(Inches(1), Inches(0.3), Inches(8), Inches(0.7))
        tf = title_box.text_frame
        tf.text = "⭐ COMPETITIVE ADVANTAGE"
        p = tf.paragraphs[0]
        p.font.size = Pt(36)
        p.font.color.rgb = COLORS['white']
        p.font.bold = True

        # Advantages
        content_box = slide.shapes.add_textbox(Inches(1), Inches(2), Inches(8), Inches(4.5))
        tf = content_box.text_frame
        tf.text = content['advantage']
        tf.word_wrap = True
        for p in tf.paragraphs:
            p.font.size = Pt(20)
            p.font.color.rgb = COLORS['dark']
            p.space_before = Pt(12)

    # ==================== 7. FINANCIALS ====================
    if content.get('financials'):
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        slide.background.fill.solid()
        slide.background.fill.fore_color.rgb = COLORS['white']

        # Header
        header = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(0), Inches(0),
            Inches(10), Inches(1.2)
        )
        header.fill.solid()
        header.fill.fore_color.rgb = COLORS['secondary']
        header.line.fill.background()

        # Title
        title_box = slide.shapes.add_textbox(Inches(1), Inches(0.3), Inches(8), Inches(0.7))
        tf = title_box.text_frame
        tf.text = "📈 FINANCIAL PROJECTIONS"
        p = tf.paragraphs[0]
        p.font.size = Pt(36)
        p.font.color.rgb = COLORS['white']
        p.font.bold = True

        # Financial content
        content_box = slide.shapes.add_textbox(Inches(1), Inches(2), Inches(8), Inches(4.5))
        tf = content_box.text_frame
        tf.text = content['financials']
        tf.word_wrap = True
        for p in tf.paragraphs:
            p.font.size = Pt(20)
            p.font.color.rgb = COLORS['dark']
            p.space_before = Pt(12)

    # ==================== 8. CALL TO ACTION ====================
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_gradient_background(slide, COLORS['accent'], COLORS['secondary'])

    # CTA box
    cta_box = slide.shapes.add_textbox(Inches(1.5), Inches(2), Inches(7), Inches(3.5))
    tf = cta_box.text_frame
    tf.text = "🚀 LET'S BUILD THE FUTURE\n\n"
    tf.text += content['cta'] + "\n\n"
    tf.text += f"Contact: {content['author']}\n"
    tf.text += f"Project: {content['project_name']}"

    for p in tf.paragraphs:
        p.alignment = PP_ALIGN.CENTER
        p.font.size = Pt(26)
        p.font.color.rgb = COLORS['white']
        p.font.bold = True
        p.space_after = Pt(20)

    # Save file
    timestamp = datetime.now().strftime('%Y%m%d_%H%M%S')
    filename = f"pitch_{package}_{user_id}_{timestamp}.pptx"
    prs.save(filename)

    logger.info(f"Created presentation: {filename}")
    return filename


# Admin approve handler'da o'zgartirish
@dp.callback_query_handler(lambda c: c.data.startswith("admin_approve"), state='*')
async def admin_approve_handler(call: types.CallbackQuery):
    """Admin tasdiqlash - yangilangan"""
    logger.info(f"Admin approval by {call.from_user.id}: {call.data}")
    await call.answer()

    if call.from_user.id != ADMIN_ID:
        await call.message.answer("⛔ Sizda ruxsat yo'q!")
        return

    user_id = int(call.data.split(":")[1])
    order = db.get_order(user_id)

    if not order:
        await call.message.answer("❌ Buyurtma topilmadi")
        return

    db.update_order(user_id, status="approved")
    db.log_admin_action(ADMIN_ID, "approve", user_id)

    await call.message.answer("⏳ To'lov tasdiqlandi. Professional PPTX tayyorlanmoqda...")

    try:
        answers = order['answers']

        # AI bilan yaxshilash
        if USE_OPENAI and openai_client:
            await call.message.answer("🤖 AI javoblarni optimizatsiya qilmoqda...")
            enhanced_answers = await enhance_answers_with_openai(answers, order['package'])
        else:
            enhanced_answers = answers

        # Mukammal PPTX yaratish
        await call.message.answer("🎨 Professional prezentatsiya yaratilmoqda...")
        pptx_path = await create_stunning_pitch_deck(user_id, enhanced_answers, order['package'])

        # Foydalanuvchiga yuborish
        package_name = 'Oddiy' if order['package'] == 'simple' else 'Professional'

        caption = (
            "🎉 Sizning professional Pitch Deck tayyor!\n\n"
            f"📦 Paket: {package_name}\n"
            f"✨ AI optimizatsiyasi: {'✅ Qo'llanildi' if USE_OPENAI else '➖ Mavjud emas'}\n"
            f"📄 Slaydlar soni: 8-10 ta\n\n"
            "🚀 Investorlarga muvaffaqiyatlar tilaymiz!\n"
            "💡 Maslahat: Prezentatsiyani ko'rib chiqing va kerak bo'lsa tahrirlang."
        )

        with open(pptx_path, "rb") as f:
            await bot.send_document(
                chat_id=user_id,
                document=InputFile(f, filename=os.path.basename(pptx_path)),
                caption=caption
            )

        db.update_order(user_id, status="completed")
        await call.message.answer(f"✅ Professional PPTX foydalanuvchi {user_id} ga yuborildi!")

        # Faylni o'chirish
        if os.path.exists(pptx_path):
            os.remove(pptx_path)
            logger.info(f"Temporary file deleted: {pptx_path}")

    except Exception as e:
        logger.error(f"PPTX generation failed: {e}")
        await call.message.answer(f"❌ Xatolik: {str(e)}")
        db.update_order(user_id, status="error")

        # Foydalanuvchiga xatolik haqida xabar
        await bot.send_message(
            chat_id=user_id,
            text="❌ Texnik xatolik yuz berdi. Admin bilan bog'laning: @support"
        )

# Qolgan kodlar o'zgarmaydi...